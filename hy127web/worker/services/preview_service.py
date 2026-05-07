import hashlib
import json
import os
from pathlib import Path

from .security import validate_path


class PreviewService:
    def __init__(self, project_root: str):
        self.project_root = project_root
        self.cache_dir = os.path.join(project_root, ".web-workbench", "preview-cache")
        os.makedirs(self.cache_dir, exist_ok=True)

    def preview(self, rel_path: str) -> dict:
        if not validate_path(self.project_root, rel_path):
            raise ValueError("路径不合法")

        abs_path = os.path.join(self.project_root, rel_path)
        if not os.path.isfile(abs_path):
            raise FileNotFoundError(f"文件不存在: {rel_path}")

        ext = os.path.splitext(rel_path)[1].lower()

        handlers = {
            ".pdf": self._preview_pdf,
            ".xlsx": self._preview_excel,
            ".docx": self._preview_word,
            ".pptx": self._preview_pptx,
        }

        handler = handlers.get(ext)
        if not handler:
            raise ValueError(f"不支持预览的文件类型: {ext}")

        return handler(abs_path, rel_path)

    def _preview_pdf(self, abs_path: str, rel_path: str) -> dict:
        return {
            "type": "pdf",
            "path": rel_path,
            "stream_url": f"/api/preview-stream?path={rel_path}",
            "note": "前端使用 PDF.js 渲染",
        }

    def get_pdf_bytes(self, rel_path: str) -> bytes:
        if not validate_path(self.project_root, rel_path):
            raise ValueError("路径不合法")
        abs_path = os.path.join(self.project_root, rel_path)
        with open(abs_path, "rb") as f:
            return f.read()

    def _preview_excel(self, abs_path: str, rel_path: str) -> dict:
        try:
            import openpyxl
        except ImportError:
            return {"type": "excel", "error": "openpyxl 未安装", "path": rel_path}

        cache_key = self._cache_key(abs_path)
        cached = self._load_cache(cache_key)
        if cached:
            return cached

        wb = openpyxl.load_workbook(abs_path, read_only=True, data_only=True)
        sheets = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(max_row=500, values_only=True):
                rows.append([self._cell_value(c) for c in row])
            sheets.append({
                "name": sheet_name,
                "rows": rows,
                "total_rows": ws.max_row or 0,
                "total_cols": ws.max_column or 0,
            })
        wb.close()

        result = {
            "type": "excel",
            "path": rel_path,
            "sheets": sheets,
            "note": "数据和基础格式准确，复杂图表/条件格式受限",
        }
        self._save_cache(cache_key, result)
        return result

    def _preview_word(self, abs_path: str, rel_path: str) -> dict:
        try:
            import mammoth
        except ImportError:
            return {"type": "word", "error": "mammoth 未安装", "path": rel_path}

        cache_key = self._cache_key(abs_path)
        cached = self._load_cache(cache_key)
        if cached:
            return cached

        with open(abs_path, "rb") as f:
            result = mammoth.convert_to_html(f)

        data = {
            "type": "word",
            "path": rel_path,
            "html": result.value,
            "warnings": [str(w) for w in result.messages[:10]],
            "note": "轻量预览：段落/表格/列表可还原，复杂排版会丢失",
        }
        self._save_cache(cache_key, data)
        return data

    def _preview_pptx(self, abs_path: str, rel_path: str) -> dict:
        try:
            from pptx import Presentation
        except ImportError:
            return {"type": "pptx", "error": "python-pptx 未安装", "path": rel_path}

        cache_key = self._cache_key(abs_path)
        cached = self._load_cache(cache_key)
        if cached:
            return cached

        prs = Presentation(abs_path)
        slides = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            slides.append({
                "index": i + 1,
                "texts": texts,
            })

        data = {
            "type": "pptx",
            "path": rel_path,
            "slides": slides,
            "total_slides": len(slides),
            "note": "内容提取模式：只能提取文本，不等于幻灯片原始排版",
        }
        self._save_cache(cache_key, data)
        return data

    def _cache_key(self, abs_path: str) -> str:
        stat = os.stat(abs_path)
        raw = f"{abs_path}:{stat.st_size}:{stat.st_mtime}"
        return hashlib.md5(raw.encode()).hexdigest()

    def _load_cache(self, key: str) -> dict | None:
        cache_file = os.path.join(self.cache_dir, f"{key}.json")
        if os.path.isfile(cache_file):
            with open(cache_file, "r", encoding="utf-8") as f:
                return json.load(f)
        return None

    def _save_cache(self, key: str, data: dict):
        cache_file = os.path.join(self.cache_dir, f"{key}.json")
        with open(cache_file, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False)

    @staticmethod
    def _cell_value(cell):
        if cell is None:
            return ""
        if isinstance(cell, (int, float)):
            return cell
        return str(cell)
